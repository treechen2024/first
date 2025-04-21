import os
import re
import argparse
from pathlib import Path
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN

def txt_to_pptx(input_file, output_file=None):
    """
    將TXT文件直接轉換為PPTX格式
    
    Args:
        input_file: 輸入的TXT文件路徑
        output_file: 輸出的PPTX文件路徑，如果為None則自動生成
    """
    if output_file is None:
        output_file = os.path.splitext(input_file)[0] + ".pptx"
    
    # 讀取TXT文件
    with open(input_file, 'r', encoding='utf-8') as f:
        content = f.read()
    
    # 分割內容為幻燈片
    # 假設每個幻燈片由至少兩個連續的換行符分隔
    slides = re.split(r'\n{2,}', content)
    
    # 創建演示文稿
    prs = Presentation()
    
    for slide in slides:
        slide = slide.strip()
        if not slide:
            continue
            
        # 處理幻燈片內容
        lines = slide.split('\n')
        
        # 假設第一行是標題
        if lines and lines[0]:
            title = lines[0]
            content_lines = lines[1:]
            
            # 創建幻燈片
            slide_layout = prs.slide_layouts[1]  # 使用標題和內容佈局
            slide = prs.slides.add_slide(slide_layout)
            
            # 設置標題
            title_shape = slide.shapes.title
            title_shape.text = title
            
            # 設置內容
            content = "\n".join(content_lines)
            content_shape = slide.placeholders[1]
            
            # 處理列表項
            text_frame = content_shape.text_frame
            text_frame.text = ""  # 清空預設文本
            
            for line in content_lines:
                line = line.strip()
                if not line:
                    continue
                
                p = text_frame.add_paragraph()
                
                # 檢測是否為列表項
                if re.match(r'^\d+\.\s', line):  # 有序列表
                    p.level = 0
                    p.text = line
                elif re.match(r'^[-*]\s', line):  # 無序列表
                    p.level = 0
                    p.text = line.replace('-', '•').replace('*', '•')
                else:
                    p.text = line
                
                # 設置字體
                for run in p.runs:
                    run.font.size = Pt(18)
    
    # 保存PPTX文件
    prs.save(output_file)
    return output_file

def main():
    parser = argparse.ArgumentParser(description='將TXT文件轉換為PPTX格式的簡報')
    parser.add_argument('input_file', help='輸入的TXT文件路徑')
    parser.add_argument('-o', '--output', help='輸出的PPTX文件路徑')
    
    args = parser.parse_args()
    
    # 轉換為PPTX格式
    pptx_file = txt_to_pptx(args.input_file, args.output)
    print(f"已生成PPTX文件: {pptx_file}")

if __name__ == "__main__":
    main()