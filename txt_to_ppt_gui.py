import tkinter as tk
from tkinter import filedialog, colorchooser, messagebox
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor


def create_pptx(txt_path, output_path, font_color, bg_color):
    prs = Presentation()
    blank_slide_layout = prs.slide_layouts[6]

    with open(txt_path, 'r', encoding='utf-8') as f:
        content = f.read()

    # 以兩個以上空白行為分頁依據
    blocks = [block.strip() for block in content.split('\n\n') if block.strip()]

    for block in blocks:
        slide = prs.slides.add_slide(blank_slide_layout)

        # 背景色
        slide.background.fill.solid()
        slide.background.fill.fore_color.rgb = RGBColor(*bg_color)

        # 新增文字方塊
        textbox = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(8), Inches(5.5))
        tf = textbox.text_frame
        tf.clear()

        # 使用 /N 作為內頁換行
        lines = block.split('/N')
        for idx, line in enumerate(lines):
            if idx == 0:
                p = tf.paragraphs[0]
            else:
                p = tf.add_paragraph()
            p.text = line.strip()
            run = p.runs[0]
            run.font.size = Pt(32)
            run.font.color.rgb = RGBColor(*font_color)

    prs.save(output_path)
    messagebox.showinfo("完成", f"PPT 產生成功：\n{output_path}")


def choose_color(label, color_var):
    color_code = colorchooser.askcolor(title="選擇顏色")
    if color_code[0]:
        rgb = tuple(int(c) for c in color_code[0])
        color_var.set(rgb)
        label.config(bg=color_code[1])


def main():
    root = tk.Tk()
    root.title("TXT 轉 PPT 工具")

    # 輸入 TXT
    tk.Label(root, text="輸入 TXT：").grid(row=0, column=0, sticky="e")
    txt_entry = tk.Entry(root, width=40)
    txt_entry.grid(row=0, column=1)
    tk.Button(root, text="選擇", command=lambda: txt_entry.insert(0, filedialog.askopenfilename(filetypes=[("Text Files", "*.txt")]))).grid(row=0, column=2)

    # 輸出 PPTX
    tk.Label(root, text="輸出 PPTX：").grid(row=1, column=0, sticky="e")
    ppt_entry = tk.Entry(root, width=40)
    ppt_entry.grid(row=1, column=1)
    tk.Button(root, text="選擇", command=lambda: ppt_entry.insert(0, filedialog.asksaveasfilename(defaultextension=".pptx", filetypes=[("PowerPoint Files", "*.pptx")]))).grid(row=1, column=2)

    # 字體顏色
    tk.Label(root, text="字體顏色：").grid(row=2, column=0, sticky="e")
    font_color_var = tk.StringVar()
    font_color_display = tk.Label(root, text="     ", bg="black")
    font_color_display.grid(row=2, column=1, sticky="w")
    tk.Button(root, text="選擇", command=lambda: choose_color(font_color_display, font_color_var)).grid(row=2, column=2)

    # 背景顏色
    tk.Label(root, text="背景顏色：").grid(row=3, column=0, sticky="e")
    bg_color_var = tk.StringVar()
    bg_color_display = tk.Label(root, text="     ", bg="white")
    bg_color_display.grid(row=3, column=1, sticky="w")
    tk.Button(root, text="選擇", command=lambda: choose_color(bg_color_display, bg_color_var)).grid(row=3, column=2)

    # 執行轉換
    def convert():
        if not txt_entry.get() or not ppt_entry.get():
            messagebox.showwarning("錯誤", "請輸入完整路徑")
            return

        # 預設顏色
        font_color = font_color_var.get() if font_color_var.get() else (0, 0, 0)
        bg_color = bg_color_var.get() if bg_color_var.get() else (255, 255, 255)

        create_pptx(
            txt_entry.get(),
            ppt_entry.get(),
            eval(font_color) if isinstance(font_color, str) else font_color,
            eval(bg_color) if isinstance(bg_color, str) else bg_color
        )

    tk.Button(root, text="轉換", width=20, command=convert).grid(row=4, column=0, columnspan=3, pady=10)

    root.mainloop()


if __name__ == "__main__":
    main()
